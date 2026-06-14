# -*- coding: utf-8 -*-
"""Surgically rewrite the code boxes on the deck, in place.

The deck (slides/*.pptx) is the master: edit it by hand in PowerPoint as much
as you like. The ONLY thing this script touches is the syntax-highlighted code
inside boxes named "code:<tag>" (created by make_deck.py's code_panel). It
re-renders each from scripts/code_snippets.py and leaves every other shape —
photos, wording, layout, the panel rectangle itself — exactly as you left it.

Workflow:
    1. edit a fragment in scripts/code_snippets.py
    2. close the deck in PowerPoint (a file lock prevents overwriting)
    3. python scripts/update_code_slides.py        # rewrites in place
    4. re-export the PDF from PowerPoint when you're happy

The code box's FONT SIZE and POSITION are preserved (read from the box), so if
you resized a box to fit, your size sticks. Only the text content + colours are
replaced. Run with --check to see what would change without writing.

Run:  python scripts/update_code_slides.py [path-to.pptx] [-o out.pptx] [--check]
"""
import argparse
import os
import re
import sys

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pygments import lex
from pygments.lexers import CppLexer

from code_snippets import SNIPPETS

# ---- theme (matches make_deck.py) -----------------------------------------
CODE_FONT    = "Consolas"
CODE_DEFAULT = RGBColor(0xD4, 0xD4, 0xD4)
YELLOW       = RGBColor(0xF2, 0xC7, 0x44)   # highlight colour
TOKEN_COLORS = {
    "Keyword":        RGBColor(0x56, 0x9C, 0xD6),
    "Keyword.Type":   RGBColor(0x4E, 0xC9, 0xB0),
    "Name.Class":     RGBColor(0x4E, 0xC9, 0xB0),
    "Name.Namespace": RGBColor(0x4E, 0xC9, 0xB0),
    "Name.Function":  RGBColor(0xDC, 0xDC, 0xAA),
    "Name.Builtin":   RGBColor(0x4E, 0xC9, 0xB0),
    "Name.Label":     RGBColor(0xDC, 0xDC, 0xAA),
    "Literal.String": RGBColor(0xCE, 0x91, 0x78),
    "Literal.Number": RGBColor(0xB5, 0xCE, 0xA8),
    "Comment":        RGBColor(0x6A, 0x99, 0x55),
    "Operator":       RGBColor(0xD4, 0xD4, 0xD4),
    "Punctuation":    RGBColor(0xD4, 0xD4, 0xD4),
}

DEFAULT_PPTX = os.path.join(
    os.path.dirname(__file__), "..", "slides",
    "ACCU-Separation-of-Concerns-at-Zero-Cost.pptx")


def token_color(ttype):
    t = str(ttype).replace("Token.", "")
    while t:
        if t in TOKEN_COLORS:
            return TOKEN_COLORS[t]
        t = t.rsplit(".", 1)[0] if "." in t else ""
    return CODE_DEFAULT


# Multi-char PascalCase identifier => a type, by convention (Component, Registry,
# LocalDerivative, TagSet, PointMass, Dual, ...). Single capitals (T, D, J, K, Q)
# are left alone — they're as often loop/temporary variables as type params.
_TYPE_NAME = re.compile(r"^[A-Z][A-Za-z0-9_]+$")


def syntax_color(ttype, text):
    """Token colour, with one fix on top of pygments: it leaves user-defined
    type names as plain grey Name tokens and is even inconsistent about it
    (e.g. 'Registry' teal after `typename`, grey in `const Registry&`). Promote
    anything that looks like a type to the type colour so types are uniform."""
    c = token_color(ttype)
    if c == CODE_DEFAULT and _TYPE_NAME.match(text):
        return TOKEN_COLORS["Keyword.Type"]
    return c


def existing_font_size(text_frame, default=Pt(14)):
    """Reuse whatever size the box currently uses, so hand-resizes stick."""
    for p in text_frame.paragraphs:
        for r in p.runs:
            if r.font.size is not None:
                return r.font.size
    return default


def render(text_frame, code, highlight_lines, size):
    """Re-render `code` into an existing text frame (clears it first).
    Mirrors make_deck.py's code_panel text rendering."""
    text_frame.clear()
    text_frame.word_wrap = False

    lines = code.rstrip("\n").split("\n")
    line_tokens = [[] for _ in lines]
    li = 0
    for ttype, value in lex(code, CppLexer()):
        parts = value.split("\n")
        for i, part in enumerate(parts):
            if i > 0:
                li += 1
            if part and li < len(line_tokens):
                line_tokens[li].append((ttype, part))

    first = True
    for idx, toks in enumerate(line_tokens):
        p = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        first = False
        p.space_after = Pt(1)
        hl = (idx + 1) in highlight_lines
        if not toks:
            r = p.add_run(); r.text = " "
            r.font.size = size; r.font.name = CODE_FONT
            continue
        for ttype, text in toks:
            r = p.add_run(); r.text = text
            r.font.size = size; r.font.name = CODE_FONT
            r.font.color.rgb = YELLOW if hl else syntax_color(ttype, text)
            if hl:
                r.font.bold = True


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("pptx", nargs="?", default=DEFAULT_PPTX,
                    help="deck to update (default: the repo deck)")
    ap.add_argument("-o", "--output",
                    help="write here instead of overwriting the input")
    ap.add_argument("--check", action="store_true",
                    help="report which code boxes were found; write nothing")
    args = ap.parse_args()

    prs = Presentation(args.pptx)

    seen = {}
    for n, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.name.startswith("code:"):
                continue
            tag = shape.name[len("code:"):]
            if tag not in SNIPPETS:
                print(f"  ! slide {n}: box 'code:{tag}' has no snippet — skipped")
                continue
            seen.setdefault(tag, []).append(n)
            if not args.check and shape.has_text_frame:
                snip = SNIPPETS[tag]
                size = existing_font_size(shape.text_frame)
                render(shape.text_frame, snip["code"], snip["highlight_lines"], size)

    for tag in sorted(seen):
        print(f"  {'found' if args.check else 'updated'} code:{tag}  (slide {','.join(map(str, seen[tag]))})")
    missing = sorted(set(SNIPPETS) - set(seen))
    if missing:
        print(f"  ! snippets with no matching box in the deck: {', '.join(missing)}")

    if args.check:
        print("check only - nothing written")
        return

    out = args.output or args.pptx
    try:
        prs.save(out)
    except PermissionError:
        sys.exit(f"ERROR: cannot write {out} — is the deck open in PowerPoint?")
    print(f"saved {out}")


if __name__ == "__main__":
    main()
